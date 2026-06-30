import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io

st.set_page_config(page_title="PPT 자동 제작 툴", layout="centered")

st.title("📊 PPT 자동 제작 툴")
st.markdown("슬라이드 내용을 입력하면 자동으로 PPT 파일을 생성합니다.")

# --- Session state ---
if "slides" not in st.session_state:
    st.session_state.slides = []

# --- Layout options ---
layout_options = {
    "제목 슬라이드": "title",
    "내용 (불릿)": "content",
    "섹션 헤더": "section",
    "이미지 + 텍스트": "image_text",
}

# --- Sidebar: Add slide ---
with st.sidebar:
    st.header("슬라이드 추가")

    layout_type = st.selectbox("레이아웃", list(layout_options.keys()))

    if layout_type == "제목 슬라이드":
        title = st.text_input("제목", key="add_title")
        subtitle = st.text_input("부제목", key="add_subtitle")

        if st.button("슬라이드 추가"):
            st.session_state.slides.append({
                "type": "title",
                "title": title,
                "subtitle": subtitle,
            })
            st.rerun()

    elif layout_type == "내용 (불릿)":
        title = st.text_input("슬라이드 제목", key="add_content_title")
        lines = st.text_area("내용 (한 줄에 하나씩)", key="add_content_lines",
                             placeholder="첫 번째 항목\n두 번째 항목\n세 번째 항목")

        if st.button("슬라이드 추가"):
            st.session_state.slides.append({
                "type": "content",
                "title": title,
                "bullets": [l.strip() for l in lines.split("\n") if l.strip()],
            })
            st.rerun()

    elif layout_type == "섹션 헤더":
        title = st.text_input("섹션 제목", key="add_section_title")

        if st.button("슬라이드 추가"):
            st.session_state.slides.append({
                "type": "section",
                "title": title,
            })
            st.rerun()

# --- Main: Current slides ---
st.subheader("슬라이드 목록")

if not st.session_state.slides:
    st.info("왼쪽 사이드바에서 슬라이드를 추가해주세요.")
else:
    for i, slide in enumerate(st.session_state.slides):
        with st.container(border=True):
            cols = st.columns([0.85, 0.15])
            with cols[0]:
                if slide["type"] == "title":
                    st.markdown(f"**#{i+1} 제목 슬라이드** — {slide['title'] or '(제목 없음)'}")
                elif slide["type"] == "content":
                    st.markdown(f"**#{i+1} 내용 슬라이드** — {slide['title'] or '(제목 없음)'} ({len(slide['bullets'])}개 항목)")
                elif slide["type"] == "section":
                    st.markdown(f"**#{i+1} 섹션 헤더** — {slide['title'] or '(제목 없음)'}")
            with cols[1]:
                if st.button("삭제", key=f"del_{i}"):
                    st.session_state.slides.pop(i)
                    st.rerun()

# --- Generate ---
if st.session_state.slides:
    if st.button("📥 PPT 생성 및 다운로드", type="primary", use_container_width=True):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_data in st.session_state.slides:
            if slide_data["type"] == "title":
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)

                txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.5))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_data["title"] or "제목 없음"
                p.font.size = Pt(44)
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER

                if slide_data.get("subtitle"):
                    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1))
                    tf2 = txBox2.text_frame
                    p2 = tf2.paragraphs[0]
                    p2.text = slide_data["subtitle"]
                    p2.font.size = Pt(24)
                    p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
                    p2.alignment = PP_ALIGN.CENTER

            elif slide_data["type"] == "content":
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)

                title_shape = slide.shapes.title
                title_shape.text = slide_data["title"] or "제목 없음"

                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()

                for idx, bullet in enumerate(slide_data["bullets"]):
                    if idx == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = bullet
                    p.font.size = Pt(20)
                    p.level = 0

            elif slide_data["type"] == "section":
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0x2E, 0x5C, 0x8F)

                txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_data["title"] or "제목 없음"
                p.font.size = Pt(40)
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        st.download_button(
            label="📎 PPT 파일 다운로드",
            data=buf,
            file_name="presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
        )

st.divider()
st.markdown("### 🚀 실행 방법")
st.code("streamlit run app.py", language="bash")
st.markdown("브라우저가 자동으로 열리며, 같은 네트워크의 다른 기기에서도 접근 가능합니다.")
st.markdown("**외부 접속**: `http://<your-ip>:8501`")
